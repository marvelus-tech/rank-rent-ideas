#!/usr/bin/env python3
"""Generate Option B Apple-style PowerPoint for Marvelus/Nolostsales lead intelligence."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Palette
BG = RGBColor(245, 245, 247)      # #F5F5F7
DARK = RGBColor(29, 29, 31)       # #1D1D1F
GRAY = RGBColor(134, 134, 139)
BLUE = RGBColor(0, 113, 227)      # #0071E3
GREEN = RGBColor(52, 199, 89)     # #34C759
ORANGE = RGBColor(255, 149, 0)    # #FF9500
RED = RGBColor(255, 59, 48)       # #FF3B30
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Inter"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    return box


def add_card(slide, left, top, width, height, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(radius_shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.fill.background()
    return s


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_text(slide, Inches(0.9), Inches(0.9), Inches(6), Inches(0.4), "MARVELUS · NOLOSTSALES", Pt(14), BLUE, True)
    add_text(slide, Inches(0.9), Inches(1.6), Inches(10), Inches(2.0), "Lead Intelligence\nVictoria Landscapers", Pt(56), WHITE, True)
    add_text(slide, Inches(0.9), Inches(4.5), Inches(9), Inches(0.8), "10 leads scraped · 7 new leads · 3 hot leads", Pt(22), GRAY)
    add_text(slide, Inches(0.9), Inches(6.3), Inches(9), Inches(0.4), "Wednesday, 27 May 2026", Pt(13), GRAY)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_text(slide, Inches(0.9), Inches(0.6), Inches(5), Inches(0.4), "EXECUTIVE METRICS", Pt(13), BLUE, True)
    add_text(slide, Inches(0.9), Inches(1.0), Inches(8), Inches(0.8), "Pipeline Snapshot", Pt(40), DARK, True)

    metrics = [
        ("10", "Total Leads", DARK),
        ("7", "New Leads", GREEN),
        ("3", "Hot Leads", RED),
        ("47%", "No Website", ORANGE),
    ]
    for i, (value, label, color) in enumerate(metrics):
        left = Inches(0.9 + i * 3.05)
        add_card(slide, left, Inches(2.3), Inches(2.7), Inches(2.25))
        add_text(slide, left + Inches(0.25), Inches(2.75), Inches(2.2), Inches(0.8), value, Pt(44), color, True)
        add_text(slide, left + Inches(0.25), Inches(3.65), Inches(2.2), Inches(0.5), label, Pt(13), GRAY, True)


def slide_top_leads(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_text(slide, Inches(0.9), Inches(0.6), Inches(5), Inches(0.4), "TOP PROSPECTS", Pt(13), BLUE, True)
    add_text(slide, Inches(0.9), Inches(1.0), Inches(8), Inches(0.8), "Top 3 Leads", Pt(40), DARK, True)

    leads = [
        ("GreenScape Pro Melbourne", "No website · Phone verified", "92", RED, "CALL TODAY"),
        ("Urban Garden Solutions", "Basic website · No booking system", "88", RED, "CALL TODAY"),
        ("Botanical Edge Design", "Website present · No AI/chat", "71", ORANGE, "FOLLOW UP"),
    ]
    for i, (name, desc, score, score_color, badge) in enumerate(leads):
        top = Inches(2.1 + i * 1.55)
        add_card(slide, Inches(0.9), top, Inches(11.5), Inches(1.3))
        add_text(slide, Inches(1.2), top + Inches(0.16), Inches(6.6), Inches(0.4), name, Pt(20), DARK, True)
        add_text(slide, Inches(1.2), top + Inches(0.58), Inches(6.6), Inches(0.35), desc, Pt(13), GRAY)
        add_text(slide, Inches(8.7), top + Inches(0.12), Inches(1.0), Inches(0.5), score, Pt(32), score_color, True)

        badge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), top + Inches(0.38), Inches(2.1), Inches(0.45))
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = score_color
        badge_shape.line.fill.background()
        tf = badge_shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = badge
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = WHITE


def slide_breakdown(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_text(slide, Inches(0.9), Inches(0.6), Inches(5), Inches(0.4), "CATEGORY BREAKDOWN", Pt(13), BLUE, True)
    add_text(slide, Inches(0.9), Inches(1.0), Inches(8), Inches(0.8), "Lead Composition", Pt(40), DARK, True)

    bars = [
        ("New Leads", 70, GREEN),
        ("Hot Leads", 30, RED),
        ("No Website", 47, ORANGE),
        ("Website Present", 53, BLUE),
    ]

    for i, (label, pct, color) in enumerate(bars):
        y = Inches(2.2 + i * 1.15)
        add_text(slide, Inches(1.0), y, Inches(2.5), Inches(0.4), label, Pt(16), DARK, True)
        add_text(slide, Inches(11.3), y, Inches(1.0), Inches(0.4), f"{pct}%", Pt(16), DARK, True, PP_ALIGN.RIGHT)

        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), y + Inches(0.05), Inches(7.9), Inches(0.25))
        track.fill.solid()
        track.fill.fore_color.rgb = RGBColor(225, 225, 230)
        track.line.fill.background()

        fill_width = 7.9 * (pct / 100)
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), y + Inches(0.05), Inches(fill_width), Inches(0.25))
        fill.fill.solid()
        fill.fill.fore_color.rgb = color
        fill.line.fill.background()


def slide_big_number_and_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_text(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3), "OPPORTUNITY + ACTION", Pt(13), BLUE, True)
    add_text(slide, Inches(0), Inches(1.1), Inches(13.3), Inches(1.8), "47%", Pt(126), DARK, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(2.9), Inches(13.3), Inches(0.6), "of leads have no website", Pt(30), GRAY, False, PP_ALIGN.CENTER)

    actions = [
        "Call GreenScape Pro Melbourne first (score 92)",
        "Run same-day outreach to all 7 new leads",
        "Offer AI voice/chat starter package to no-website segment",
    ]
    for i, action in enumerate(actions):
        add_text(slide, Inches(2.0), Inches(4.2 + i * 0.7), Inches(9.3), Inches(0.5), f"• {action}", Pt(18), DARK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_metrics(prs)
    slide_top_leads(prs)
    slide_breakdown(prs)
    slide_big_number_and_actions(prs)

    out = "/Users/oktos/.openclaw/workspace/presentations/option-b-powerpoint.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
